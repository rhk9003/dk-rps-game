"""
單頁簡報：機率與數量分配（討論用主表）
DK 府中店滿月｜閨密機猜拳遊戲
排版：贏家上、輸家下；字放大
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CORAL = RGBColor(0xF9, 0x61, 0x67)
GOLD  = RGBColor(0xBF, 0x8F, 0x00)
NAVY  = RGBColor(0x2F, 0x3C, 0x7E)
CREAM = RGBColor(0xFC, 0xF6, 0xF5)
INK   = RGBColor(0x24, 0x24, 0x24)
MUTED = RGBColor(0x77, 0x77, 0x77)
WIN_BG  = RGBColor(0xFF, 0xF2, 0xCC)
LOSE_BG = RGBColor(0xDE, 0xEB, 0xF7)
WIN_BG_ALT  = RGBColor(0xFF, 0xF9, 0xDE)
LOSE_BG_ALT = RGBColor(0xEE, 0xF4, 0xFC)
WIN_ACCENT  = RGBColor(0xF4, 0xB0, 0x84)
LOSE_ACCENT = RGBColor(0x9D, 0xC3, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# 背景
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()
bg.shadow.inherit = False

# 標題條
HEADER_H = Inches(0.85)
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
header.fill.solid()
header.fill.fore_color.rgb = INK
header.line.fill.background()

title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.08), Inches(12.5), Inches(0.55))
title_tf = title_tb.text_frame
title_tf.margin_left = title_tf.margin_right = 0
title_tf.margin_top = title_tf.margin_bottom = 0
p = title_tf.paragraphs[0]
p.text = '📌 機率與數量分配（討論用主表）'
p.font.name = 'Calibri'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = WHITE

sub_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.3))
sub_tf = sub_tb.text_frame
sub_tf.margin_left = sub_tf.margin_right = 0
sub_tf.margin_top = sub_tf.margin_bottom = 0
sp = sub_tf.paragraphs[0]
sp.text = 'DK 府中店滿月｜閨密機猜拳遊戲    ·    30 天 × 25 人/日 = 750 場    ·    贏率 50%    ·    券走 LINE bot 不限張數'
sp.font.name = 'Calibri'
sp.font.size = Pt(13)
sp.font.color.rgb = RGBColor(0xD5, 0xD5, 0xD5)

# ====================================================================
# 共用
# ====================================================================
def add_rect(left, top, width, height, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_text(left, top, width, height, text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font='Calibri', anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb

# 欄寬：總寬 12.5 寸
LEFT_X = Inches(0.42)
COL_W  = Inches(12.5)
col_widths = [Inches(4.5), Inches(1.8), Inches(1.5), Inches(1.8), Inches(2.9)]
col_xs = [LEFT_X]
for w in col_widths[:-1]:
    col_xs.append(col_xs[-1] + w)
headers = ['項目', '單價/面額', '機率', '數量上限', '每日上限']

BANNER_H = Inches(0.45)
HDR_H    = Inches(0.35)
ROW_H    = Inches(0.36)
SUB_H    = Inches(0.38)

def draw_pool(top, banner_text, banner_color, rows, sub_text, sub_color, row_bg, row_bg_alt):
    y = top
    # Banner
    add_rect(LEFT_X, y, COL_W, BANNER_H, banner_color)
    add_text(LEFT_X, y, COL_W, BANNER_H, banner_text,
             size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y += BANNER_H

    # Header
    add_rect(LEFT_X, y, COL_W, HDR_H, INK)
    for h, x, w in zip(headers, col_xs, col_widths):
        add_text(x, y, w, HDR_H, h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y += HDR_H

    # Rows
    for ri, row_data in enumerate(rows):
        fill = row_bg if ri % 2 == 0 else row_bg_alt
        add_rect(LEFT_X, y, COL_W, ROW_H, fill, line_rgb=RGBColor(0xDC, 0xDC, 0xDC))
        for x, w, val, i in zip(col_xs, col_widths, row_data, range(len(row_data))):
            bold = (i == 0)
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
            add_text(x, y, w, ROW_H, val, size=14, bold=bold, color=INK, align=align)
        y += ROW_H

    # Subtotal
    add_rect(LEFT_X, y, COL_W, SUB_H, sub_color)
    add_text(LEFT_X + Inches(0.2), y, COL_W - Inches(0.4), SUB_H,
             sub_text, size=13, bold=True, color=INK, align=PP_ALIGN.LEFT)
    return y + SUB_H

# ====================================================================
# 🏆 贏家獎（上）
# ====================================================================
win_rows = [
    ('A0220 DK 動態釋壓鞋',   '$1,280', '8%',   '30 件', '每日 1 雙'),
    ('A0119 五趾隱形襪',       '$250',   '25%',  '60 件', '每日 2 雙'),
    ('LINE $500 官網抵用金',   '$500',   '20%',  '不限',   '不限'),
    ('LINE 門市 7 折券',       '7 折',   '25%',  '不限',   '不限'),
    ('LINE 門市 $500 折價券',  '$500',   '22%',  '不限',   '不限'),
]

win_end = draw_pool(
    top=Inches(0.95),
    banner_text='🏆  贏家獎池    |    機率合計 100%    |    贏家場次 375 場',
    banner_color=GOLD,
    rows=win_rows,
    sub_text='贏家池 100%    ｜    實體上限 90 件（A0220×30 + A0119×60）    ｜    LINE 券不限張數',
    sub_color=WIN_ACCENT,
    row_bg=WIN_BG,
    row_bg_alt=WIN_BG_ALT,
)

# ====================================================================
# 🎁 輸家獎（下）
# ====================================================================
lose_rows = [
    ('A0003 石墨烯袖套 🎊 隱藏', '$680', '2.5%',  '10 件', '抽完為止'),
    ('A0114 涼感底紗襪',         '$180', '20%',   '60 件', '每日 2 雙'),
    ('LINE $50 官網抵用金',      '$50',  '20%',   '不限',   '不限'),
    ('LINE 門市 8 折券',          '8 折', '27.5%', '不限',   '不限'),
    ('LINE 門市 $100 折價券',    '$100', '30%',   '不限',   '不限'),
]

lose_end = draw_pool(
    top=win_end + Inches(0.1),
    banner_text='🎁  輸家獎 / 參加獎池    |    機率合計 100%    |    輸家場次 375 場',
    banner_color=NAVY,
    rows=lose_rows,
    sub_text='輸家池 100%    ｜    實體上限 70 件（A0003×10 + A0114×60）    ｜    LINE 券不限張數',
    sub_color=LOSE_ACCENT,
    row_bg=LOSE_BG,
    row_bg_alt=LOSE_BG_ALT,
)

# ====================================================================
# 底部備註
# ====================================================================
note_y = lose_end + Inches(0.05)
note_text = ('門市券：限原價商品 · ~5/31 · 不得併用      '
             'A0003（稀有隱藏）抽完自動改發 A0114      '
             '實體優先用府中現貨，不足再調      '
             '實體成本上限 ≈ $71,000  /  LINE 券邊際成本 ≈ 0')
add_text(LEFT_X, note_y, COL_W, Inches(0.3),
         note_text, size=10.5, color=MUTED, align=PP_ALIGN.LEFT)

# ====================================================================
# ============================= 第 2 頁 ==============================
# 優惠券使用規則
# ====================================================================
slide2 = prs.slides.add_slide(blank)

bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg2.fill.solid()
bg2.fill.fore_color.rgb = CREAM
bg2.line.fill.background()
bg2.shadow.inherit = False

header2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
header2.fill.solid()
header2.fill.fore_color.rgb = INK
header2.line.fill.background()

# 第二頁共用函式
def add_rect2(left, top, width, height, fill_rgb, line_rgb=None):
    shp = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_text2(left, top, width, height, text, size=14, bold=False, color=INK,
              align=PP_ALIGN.LEFT, font='Calibri', anchor=MSO_ANCHOR.MIDDLE):
    tb = slide2.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb

# 標題
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.08), Inches(12.5), Inches(0.55))
tf2 = title2.text_frame
tf2.margin_left = tf2.margin_right = 0
tf2.margin_top = tf2.margin_bottom = 0
p2 = tf2.paragraphs[0]
p2.text = '🎫 優惠券使用規則'
p2.font.name = 'Calibri'
p2.font.size = Pt(30)
p2.font.bold = True
p2.font.color.rgb = WHITE

sub2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.3))
sub2.text_frame.margin_left = sub2.text_frame.margin_right = 0
sub2.text_frame.margin_top = sub2.text_frame.margin_bottom = 0
ss = sub2.text_frame.paragraphs[0]
ss.text = 'DK 府中店滿月｜閨密機猜拳遊戲    ·    所有券透過 LINE bot 自動發送、系統核銷'
ss.font.name = 'Calibri'
ss.font.size = Pt(13)
ss.font.color.rgb = RGBColor(0xD5, 0xD5, 0xD5)

# ====================================================================
# 左欄：💻 電商券
# ====================================================================
EC_COLOR   = RGBColor(0xE8, 0x7A, 0x35)   # 暖橘
EC_BG      = RGBColor(0xFD, 0xED, 0xDC)
EC_BG_ALT  = RGBColor(0xFC, 0xE4, 0xCC)
ST_COLOR   = RGBColor(0x3E, 0x8E, 0x5A)   # 綠
ST_BG      = RGBColor(0xE2, 0xEF, 0xDA)
ST_BG_ALT  = RGBColor(0xD0, 0xE6, 0xC1)

LEFT_X2   = Inches(0.42)
LEFT_W    = Inches(6.25)
RIGHT_X2  = Inches(6.83)
RIGHT_W   = Inches(6.1)
COL_TOP   = Inches(0.95)

# 左：電商券標題條
add_rect2(LEFT_X2, COL_TOP, LEFT_W, Inches(0.55), EC_COLOR)
add_text2(LEFT_X2, COL_TOP, LEFT_W, Inches(0.55),
          '💻  電商輪盤券（官網 / APP 使用）',
          size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 電商券卡片
ec_cards = [
    {
        'pool': '🎁 輸家',
        'pool_color': LOSE_ACCENT,
        'name': 'DK 官網 $50 抵用金',
        'face': '$50',
        'scope': 'DK 官網 / APP 購買鞋款、包類商品',
        'threshold': '單筆滿 $51 以上可折抵',
        'period': '即日起 ~ 2026/5/31',
    },
    {
        'pool': '🏆 贏家',
        'pool_color': WIN_ACCENT,
        'name': 'DK 官網 $500 抵用金',
        'face': '$500',
        'scope': 'DK 官網 / APP 購買原價鞋款（不含拖鞋）',
        'threshold': '單筆滿 $1,800 以上可折抵',
        'period': '即日起 ~ 2026/5/31',
    },
]

y = COL_TOP + Inches(0.65)
card_h = Inches(1.9)
for i, cp in enumerate(ec_cards):
    bg_fill = EC_BG if i % 2 == 0 else EC_BG_ALT
    add_rect2(LEFT_X2, y, LEFT_W, card_h, bg_fill, line_rgb=EC_COLOR)
    # 池標籤
    add_rect2(LEFT_X2 + Inches(0.2), y + Inches(0.15), Inches(0.9), Inches(0.3), cp['pool_color'])
    add_text2(LEFT_X2 + Inches(0.2), y + Inches(0.15), Inches(0.9), Inches(0.3),
              cp['pool'], size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    # 面額 (大)
    add_text2(LEFT_X2 + LEFT_W - Inches(2.0), y + Inches(0.1), Inches(1.8), Inches(0.55),
              cp['face'], size=32, bold=True, color=EC_COLOR, align=PP_ALIGN.RIGHT)
    # 名稱
    add_text2(LEFT_X2 + Inches(0.2), y + Inches(0.5), LEFT_W - Inches(0.4), Inches(0.4),
              cp['name'], size=16, bold=True, color=INK, align=PP_ALIGN.LEFT)
    # 詳細 (3行)
    detail_y = y + Inches(0.95)
    for idx, (label, val) in enumerate([('適用', cp['scope']), ('門檻', cp['threshold']), ('期限', cp['period'])]):
        dy = detail_y + Inches(0.3 * idx)
        add_text2(LEFT_X2 + Inches(0.25), dy, Inches(0.6), Inches(0.28),
                  label, size=12, bold=True, color=EC_COLOR, align=PP_ALIGN.LEFT)
        add_text2(LEFT_X2 + Inches(0.85), dy, LEFT_W - Inches(1.05), Inches(0.28),
                  val, size=12, color=INK, align=PP_ALIGN.LEFT)
    y += card_h + Inches(0.1)

# ====================================================================
# 右欄：🏪 門市券
# ====================================================================
add_rect2(RIGHT_X2, COL_TOP, RIGHT_W, Inches(0.55), ST_COLOR)
add_text2(RIGHT_X2, COL_TOP, RIGHT_W, Inches(0.55),
          '🏪  門市券（DK 實體店使用）',
          size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 門市券通用規則框
rules_y = COL_TOP + Inches(0.65)
add_rect2(RIGHT_X2, rules_y, RIGHT_W, Inches(0.7), WHITE, line_rgb=ST_COLOR)
add_text2(RIGHT_X2 + Inches(0.2), rules_y, Inches(1.1), Inches(0.7),
          '通用規則', size=13, bold=True, color=ST_COLOR, align=PP_ALIGN.LEFT)
add_text2(RIGHT_X2 + Inches(1.3), rules_y, RIGHT_W - Inches(1.5), Inches(0.7),
          '✓ 限原價商品使用    ✓ 即日起 ~ 2026/5/31    ✓ 不得與其他優惠併用',
          size=12.5, color=INK, align=PP_ALIGN.LEFT)

# 4 張門市券表格
store_coupons = [
    ('🏆 贏家', WIN_ACCENT, '門市 7 折券',   '7 折', '贏家池最大折扣'),
    ('🏆 贏家', WIN_ACCENT, '門市 $500 折價券', '$500', '贏家池現金折抵'),
    ('🎁 輸家', LOSE_ACCENT, '門市 8 折券',   '8 折', '輸家池大宗 27.5%'),
    ('🎁 輸家', LOSE_ACCENT, '門市 $100 折價券', '$100', '輸家池最高機率 30%'),
]

y = rules_y + Inches(0.8)
# 表頭
hdr_h = Inches(0.38)
add_rect2(RIGHT_X2, y, RIGHT_W, hdr_h, ST_COLOR)
sc_widths = [Inches(1.0), Inches(2.3), Inches(1.3), Inches(1.5)]
sc_xs = [RIGHT_X2]
for w in sc_widths[:-1]:
    sc_xs.append(sc_xs[-1] + w)
sc_hdrs = ['獎池', '券名稱', '面額/折扣', '備註']
for h, x, w in zip(sc_hdrs, sc_xs, sc_widths):
    add_text2(x, y, w, hdr_h, h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
y += hdr_h

# 4 列
row_h2 = Inches(0.5)
for i, (pool, pcolor, name, face, note) in enumerate(store_coupons):
    fill = ST_BG if i % 2 == 0 else ST_BG_ALT
    add_rect2(RIGHT_X2, y, RIGHT_W, row_h2, fill, line_rgb=RGBColor(0xCC, 0xCC, 0xCC))
    # 池
    pc = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 sc_xs[0] + Inches(0.1), y + Inches(0.1),
                                 sc_widths[0] - Inches(0.2), row_h2 - Inches(0.2))
    pc.fill.solid()
    pc.fill.fore_color.rgb = pcolor
    pc.line.fill.background()
    pc.shadow.inherit = False
    add_text2(sc_xs[0], y, sc_widths[0], row_h2,
              pool, size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    # 名稱
    add_text2(sc_xs[1], y, sc_widths[1], row_h2,
              name, size=14, bold=True, color=INK, align=PP_ALIGN.LEFT)
    # 面額
    add_text2(sc_xs[2], y, sc_widths[2], row_h2,
              face, size=18, bold=True, color=ST_COLOR, align=PP_ALIGN.CENTER)
    # 備註
    add_text2(sc_xs[3], y, sc_widths[3], row_h2,
              note, size=11, color=MUTED, align=PP_ALIGN.LEFT)
    y += row_h2

# ====================================================================
# 底部共同條款
# ====================================================================
NOTE_Y = Inches(6.55)
add_rect2(Inches(0.42), NOTE_Y, Inches(12.5), Inches(0.6), WHITE, line_rgb=CORAL)
add_text2(Inches(0.55), NOTE_Y, Inches(2.5), Inches(0.6),
          '📋 共同規定',
          size=14, bold=True, color=CORAL, align=PP_ALIGN.LEFT)
add_text2(Inches(2.8), NOTE_Y, Inches(10.0), Inches(0.6),
          '• LINE bot 發送、每個 LINE ID 綁一張、逾期自動失效    '
          '• 逾期無效、不得折抵現金、不得退換    '
          '• DK 高博士保留修改及解釋活動之權利',
          size=11.5, color=INK, align=PP_ALIGN.LEFT)

# 底部來源
src2 = slide2.shapes.add_textbox(Inches(0.4), Inches(7.25), Inches(12.5), Inches(0.2))
src2.text_frame.margin_top = 0
src2.text_frame.margin_bottom = 0
sr = src2.text_frame.paragraphs[0]
sr.text = '資料來源：府中店滿月_輪盤禮.docx'
sr.font.name = 'Calibri'
sr.font.size = Pt(8)
sr.font.color.rgb = MUTED
sr.alignment = PP_ALIGN.RIGHT

out_path = '../府中店滿月_機率與數量分配_簡報.pptx'
prs.save(out_path)
import os
print('saved to:', os.path.abspath(out_path))
